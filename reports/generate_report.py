"""
Rushtown Poultry — Automated Report Generator
Pulls from Firebase, generates Excel + PowerPoint, emails via Resend.
Usage: python generate_report.py --type daily|weekly|monthly
"""

import os
import sys
import json
import argparse
import base64
import io
from datetime import datetime, timedelta, date
from calendar import monthrange

import firebase_admin
from firebase_admin import credentials, firestore

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import resend

# ── Constants ──────────────────────────────────────────────────────────────────
SEND_TO   = "jschmidt@rushtownpoultry.com"
SEND_FROM = "jbschmidt827@gmail.com"
FARMS     = {"Hegins": 8, "Danville": 5}

# EOS scope: 13 maintenance-tracked houses (Hegins H1-H8 + Danville H1-H5)
EOS_HOUSES = (
    [("Hegins",   str(i)) for i in range(1, 9)] +
    [("Danville", str(i)) for i in range(1, 6)]
)

GREEN_DARK  = RGBColor(0x1a, 0x3a, 0x1a)
GREEN_MID   = RGBColor(0x2a, 0x6a, 0x2a)
GREEN_LIGHT = RGBColor(0x4c, 0xaf, 0x50)
WHITE       = RGBColor(0xff, 0xff, 0xff)
RED         = RGBColor(0xe5, 0x3e, 0x3e)
AMBER       = RGBColor(0xd6, 0x9e, 0x2e)
GREY        = RGBColor(0x5a, 0x5a, 0x5a)

# ── Firebase ───────────────────────────────────────────────────────────────────
def init_firebase():
    sa_json = os.environ.get("FIREBASE_SERVICE_ACCOUNT")
    if not sa_json:
        sys.exit("Missing FIREBASE_SERVICE_ACCOUNT env var")
    cred = credentials.Certificate(json.loads(sa_json))
    firebase_admin.initialize_app(cred)
    return firestore.client()


def fetch_walks(db, start_date: date, end_date: date):
    """Return barnWalks records in [start_date, end_date] inclusive."""
    docs = (db.collection("barnWalks")
              .where("date", ">=", str(start_date))
              .where("date", "<=", str(end_date))
              .stream())
    return [d.to_dict() for d in docs]


def fetch_morning_walks(db, start_date: date, end_date: date):
    docs = (db.collection("morningWalks")
              .where("date", ">=", str(start_date))
              .where("date", "<=", str(end_date))
              .stream())
    return [d.to_dict() for d in docs]


def fetch_work_orders(db, start_date: date, end_date: date):
    docs = (db.collection("workOrders")
              .where("date", ">=", str(start_date))
              .where("date", "<=", str(end_date))
              .stream())
    return [d.to_dict() for d in docs]



# ── EOS snapshot helpers ───────────────────────────────────────────────────────
def fetch_open_work_orders(db, limit=500):
    """Return open work orders using server-side status filter.
    Status values in the data: 'open' | 'in-progress' | 'on-hold' | 'completed'."""
    OPEN_STATUSES = ["open", "in-progress", "on-hold"]
    out = []
    try:
        q = (db.collection("workOrders")
               .where("status", "in", OPEN_STATUSES)
               .limit(limit))
        for d in q.stream():
            wo = d.to_dict() or {}
            wo["_id"] = d.id
            out.append(wo)
        return out
    except Exception as e:
        print(f"WARN: server-side WO filter failed: {e}", file=sys.stderr)
    # Fallback: scan & filter client-side
    closed = {"completed", "complete", "closed", "resolved", "cancelled"}
    for d in db.collection("workOrders").limit(limit * 4).stream():
        wo = d.to_dict() or {}
        wo["_id"] = d.id
        st = str(wo.get("status", "")).strip().lower()
        if st and st not in closed:
            out.append(wo)
        if len(out) >= limit:
            break
    return out


def fetch_active_red_tags(db, limit=500):
    """Return redTags whose status is 'Tagged' or 'Under Review', case-insensitive."""
    active_lower = {"tagged", "under review"}
    out = []
    for d in db.collection("redTags").limit(limit).stream():
        rt = d.to_dict() or {}
        rt["_id"] = d.id
        if str(rt.get("status", "")).strip().lower() in active_lower:
            out.append(rt)
    return out


def _norm_house(h):
    """Extract just the house number from values like '5', 'H5', 'House 5', 'house 5'.
    Returns '' for 'N/A', 'PM-Generated', or anything without a digit."""
    if h is None:
        return ""
    import re as _re
    m = _re.search(r"\d+", str(h))
    return m.group(0) if m else ""


def build_eos_snapshot(db, today):
    """Build the EOS daily snapshot dict — written to data/eos-snapshot.json."""
    open_wos    = fetch_open_work_orders(db)
    red_tags    = fetch_active_red_tags(db)
    today_str   = str(today)
    today_walks = fetch_walks(db, today, today)
    today_mws   = fetch_morning_walks(db, today, today)

    now_ms = int(datetime.utcnow().timestamp() * 1000)
    seven_days_ms = 7 * 24 * 3600 * 1000

    by_house = {}
    for farm, house in EOS_HOUSES:
        by_house[f"{farm}|{house}"] = {
            "farm": farm, "house": house,
            "openWO": 0, "redTags": 0, "overdue7d": 0,
            "walkedToday": False, "morningWalkToday": False,
            "topItem": None, "topItemDays": None,
        }

    def bucket(farm, house):
        return by_house.get(f"{farm}|{_norm_house(house)}")

    for wo in open_wos:
        b = bucket(wo.get("farm",""), wo.get("house",""))
        if b:
            b["openWO"] += 1
            ts = wo.get("ts") or wo.get("createdAt") or 0
            try: ts = int(ts)
            except Exception: ts = 0
            age_days = (now_ms - ts) // (24*3600*1000) if ts else None
            if not b["topItem"] or (age_days and (b["topItemDays"] or 0) < age_days):
                b["topItem"] = wo.get("title") or wo.get("problem") or (wo.get("description","") or "")[:60]
                b["topItemDays"] = age_days

    for rt in red_tags:
        b = bucket(rt.get("farm",""), rt.get("house",""))
        if b:
            b["redTags"] += 1
            ts = rt.get("ts") or rt.get("createdAt") or 0
            try: ts = int(ts)
            except Exception: ts = 0
            if ts and (now_ms - ts) > seven_days_ms:
                b["overdue7d"] += 1

    for w in today_walks:
        b = bucket(w.get("farm",""), w.get("house",""))
        if b: b["walkedToday"] = True
    for w in today_mws:
        b = bucket(w.get("farm",""), w.get("house",""))
        if b: b["morningWalkToday"] = True

    headline = {
        "openWOTotal":   sum(b["openWO"]    for b in by_house.values()),
        "redTagsActive": sum(b["redTags"]   for b in by_house.values()),
        "overdue7d":     sum(b["overdue7d"] for b in by_house.values()),
        "walksToday":    sum(1 for b in by_house.values() if b["walkedToday"] or b["morningWalkToday"]),
        "totalHouses":   len(EOS_HOUSES),
    }

    def slim_wo(wo):
        return {
            "id":       wo.get("_id"),
            "farm":     wo.get("farm",""),
            "house":    _norm_house(wo.get("house","")),
            "title":    wo.get("title") or wo.get("problem") or "",
            "desc":     (wo.get("description") or "")[:200],
            "priority": wo.get("priority",""),
            "status":   wo.get("status",""),
            "ts":       wo.get("ts") or wo.get("createdAt"),
            "urgent":   bool(wo.get("urgent")),
        }

    def slim_rt(rt):
        ts = rt.get("ts") or rt.get("createdAt") or 0
        try: ts = int(ts)
        except Exception: ts = 0
        days_old = ((now_ms - ts) // (24*3600*1000)) if ts else None
        return {
            "id":      rt.get("_id"),
            "farm":    rt.get("farm",""),
            "house":   _norm_house(rt.get("house","")),
            "item":    rt.get("item") or rt.get("description","") or "",
            "status":  rt.get("status",""),
            "daysOld": days_old,
            "ts":      ts,
        }

    return {
        "generatedAt": datetime.utcnow().isoformat() + "Z",
        "date":        today_str,
        "scope":       {"farms": list({h[0] for h in EOS_HOUSES}), "houses": len(EOS_HOUSES)},
        "headline":    headline,
        "byHouse":     list(by_house.values()),
        "openWOs":     [slim_wo(w) for w in open_wos],
        "redTags":     [slim_rt(r) for r in red_tags],
        "walksToday":  [{"farm": w.get("farm",""), "house": _norm_house(w.get("house","")),
                         "employee": w.get("employee","") or w.get("name",""),
                         "time": w.get("time","")} for w in today_walks],
        "morningWalksToday": [{"farm": w.get("farm",""), "house": _norm_house(w.get("house","")),
                                "employee": w.get("employee","") or w.get("name",""),
                                "time": w.get("time","")} for w in today_mws],
    }


def write_eos_snapshot(snapshot, path="data/eos-snapshot.json"):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(snapshot, f, indent=2, default=str)
    print(f"Wrote EOS snapshot -> {path} "
          f"({snapshot['headline']['openWOTotal']} open WOs, "
          f"{snapshot['headline']['redTagsActive']} red tags, "
          f"{snapshot['headline']['walksToday']}/{snapshot['headline']['totalHouses']} walks today)")


# ── Date helpers ───────────────────────────────────────────────────────────────
def date_range(report_type: str):
    today = date.today()
    if report_type == "daily":
        d = today - timedelta(days=1)
        return d, d
    elif report_type == "weekly":
        end   = today - timedelta(days=1)
        start = end - timedelta(days=6)
        return start, end
    else:  # monthly
        first = today.replace(day=1) - timedelta(days=1)
        start = first.replace(day=1)
        return start, first


# ── Excel ──────────────────────────────────────────────────────────────────────
HDR_FILL  = PatternFill("solid", fgColor="1A3A1A")
ALT_FILL  = PatternFill("solid", fgColor="F4FAF4")
FLAG_FILL = PatternFill("solid", fgColor="FFF0F0")
HDR_FONT  = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
BODY_FONT = Font(name="Calibri", size=10)
FLAG_FONT = Font(name="Calibri", size=10, color="CC0000")
THIN      = Side(style="thin", color="CCCCCC")
BORDER    = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)


def _hdr(ws, row, col, text, width=None):
    c = ws.cell(row=row, column=col, value=text)
    c.font = HDR_FONT
    c.fill = HDR_FILL
    c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    c.border = BORDER
    if width:
        ws.column_dimensions[get_column_letter(col)].width = width


def _cell(ws, row, col, value, flagged=False, center=False):
    c = ws.cell(row=row, column=col, value=value)
    c.font = FLAG_FONT if flagged else BODY_FONT
    c.fill = FLAG_FILL if flagged else (ALT_FILL if row % 2 == 0 else PatternFill())
    c.alignment = Alignment(horizontal="center" if center else "left", vertical="center")
    c.border = BORDER


def build_excel(walks, m_walks, report_type: str, start_date: date, end_date: date):
    wb = openpyxl.Workbook()

    # ── Sheet 1: Summary ──────────────────────────────────────────────────────
    ws = wb.active
    ws.title = "Summary"
    ws.row_dimensions[1].height = 30
    period = f"{start_date}" if start_date == end_date else f"{start_date} to {end_date}"
    ws["A1"] = f"Rushtown Poultry — {report_type.title()} Production Report"
    ws["A1"].font = Font(bold=True, size=14, name="Calibri", color="1A3A1A")
    ws["A2"] = f"Period: {period}"
    ws["A2"].font = Font(size=11, name="Calibri", color="5A5A5A")

    kpis = [
        ("Total Barn Checks", len(walks)),
        ("Total Mortality", sum(int(w.get("mortCount") or 0) for w in walks)),
        ("Loose Birds Total", sum(int(w.get("looseCount") or 0) for w in walks)),
        ("Flagged Checks", sum(1 for w in walks if w.get("flags"))),
        ("Morning Walks", len(m_walks)),
        ("Avg Water PSI", round(sum(float(w.get("waterPSI") or 0) for w in m_walks) / max(len(m_walks),1), 1)),
        ("Avg House Temp (°F)", round(sum(float(w.get("temp") or 0) for w in m_walks) / max(len(m_walks),1), 1)),
    ]
    ws.cell(row=4, column=1, value="KPI").font = HDR_FONT
    ws.cell(row=4, column=1).fill = HDR_FILL
    ws.cell(row=4, column=1).border = BORDER
    ws.cell(row=4, column=2, value="Value").font = HDR_FONT
    ws.cell(row=4, column=2).fill = HDR_FILL
    ws.cell(row=4, column=2).border = BORDER
    ws.column_dimensions["A"].width = 28
    ws.column_dimensions["B"].width = 16
    for i, (k, v) in enumerate(kpis, start=5):
        _cell(ws, i, 1, k)
        _cell(ws, i, 2, v, center=True)

    # ── Sheet 2: Barn Checks ──────────────────────────────────────────────────
    ws2 = wb.create_sheet("Barn Checks")
    headers = ["Date","Farm","Barn","Employee","Time","Mortality","Loose Birds",
               "Water PSI","House Temp","Air Quality","Feeders","Feathering",
               "Egg Belts","Manure Belts","Rodents","Flags"]
    widths  = [12, 10, 6, 18, 8, 10, 12, 10, 11, 12, 10, 12, 10, 13, 8, 40]
    for i, (h, w) in enumerate(zip(headers, widths), 1):
        _hdr(ws2, 1, i, h, w)
    ws2.row_dimensions[1].height = 20

    for r, w in enumerate(sorted(walks, key=lambda x: (x.get("date",""), x.get("farm",""), int(x.get("house",0)))), 2):
        flagged = bool(w.get("flags"))
        _cell(ws2, r, 1,  w.get("date",""),         flagged)
        _cell(ws2, r, 2,  w.get("farm",""),         flagged)
        _cell(ws2, r, 3,  w.get("house",""),        flagged, center=True)
        _cell(ws2, r, 4,  w.get("employee",""),     flagged)
        _cell(ws2, r, 5,  w.get("time",""),         flagged, center=True)
        _cell(ws2, r, 6,  w.get("mortCount") or 0, flagged, center=True)
        _cell(ws2, r, 7,  w.get("looseCount") or 0,flagged, center=True)
        _cell(ws2, r, 8,  w.get("waterPSI") or "", flagged, center=True)
        _cell(ws2, r, 9,  w.get("temp") or "",     flagged, center=True)
        _cell(ws2, r, 10, w.get("air","").upper() if w.get("air") else "", flagged, center=True)
        _cell(ws2, r, 11, w.get("feed","").upper() if w.get("feed") else "", flagged, center=True)
        _cell(ws2, r, 12, w.get("feather","").upper() if w.get("feather") else "", flagged, center=True)
        _cell(ws2, r, 13, w.get("eggbelt","").upper() if w.get("eggbelt") else "", flagged, center=True)
        _cell(ws2, r, 14, w.get("manure","").upper() if w.get("manure") else "", flagged, center=True)
        _cell(ws2, r, 15, "YES" if w.get("rodent")=="yes" else "no", flagged, center=True)
        _cell(ws2, r, 16, "; ".join(w.get("flags", [])), flagged)

    # ── Sheet 3: Morning Walks ─────────────────────────────────────────────────
    ws3 = wb.create_sheet("Morning Walks")
    mh = ["Date","Farm","Barn","Lead / WNO","Time","EE Count","Water PSI","House Temp","Feed","Fans","Blowers","Flags"]
    mw = [12, 10, 6, 18, 8, 10, 10, 12, 8, 8, 10, 40]
    for i, (h, w) in enumerate(zip(mh, mw), 1):
        _hdr(ws3, 1, i, h, w)
    for r, mwk in enumerate(sorted(m_walks, key=lambda x: (x.get("date",""), x.get("farm",""), int(x.get("house",0)))), 2):
        flagged = bool(mwk.get("flags"))
        _cell(ws3, r, 1,  mwk.get("date",""),         flagged)
        _cell(ws3, r, 2,  mwk.get("farm",""),         flagged)
        _cell(ws3, r, 3,  mwk.get("house",""),        flagged, center=True)
        _cell(ws3, r, 4,  mwk.get("employee",""),     flagged)
        _cell(ws3, r, 5,  mwk.get("time",""),         flagged, center=True)
        _cell(ws3, r, 6,  mwk.get("eeCount") or "",  flagged, center=True)
        _cell(ws3, r, 7,  mwk.get("waterPSI") or "", flagged, center=True)
        _cell(ws3, r, 8,  mwk.get("temp") or "",     flagged, center=True)
        _cell(ws3, r, 9,  mwk.get("feed","").upper() if mwk.get("feed") else "", flagged, center=True)
        _cell(ws3, r, 10, mwk.get("fans","").upper() if mwk.get("fans") else "", flagged, center=True)
        _cell(ws3, r, 11, mwk.get("blowers","").upper() if mwk.get("blowers") else "", flagged, center=True)
        _cell(ws3, r, 12, "; ".join(mwk.get("flags", [])), flagged)

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ── PowerPoint ─────────────────────────────────────────────────────────────────
def _add_slide(prs, layout_idx=6):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _txbox(slide, left, top, width, height, text, size=18, bold=False,
           color=None, align=PP_ALIGN.LEFT, bg=None):
    from pptx.util import Emu
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    if color:
        run.font.color.rgb = color
    if bg:
        from pptx.oxml.ns import qn
        from lxml import etree
        sp = txb._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(sp, qn("p:spPr"))
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgbClr.set("val", bg)
    return txb


def _rect(slide, left, top, width, height, fill_hex):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    shape.line.fill.background()
    return shape


def _kpi_box(slide, left, top, label, value, color_hex="4caf50"):
    _rect(slide, left, top, 2.1, 1.1, "1a3a1a")
    _txbox(slide, left+0.1, top+0.05, 1.9, 0.55, str(value),
           size=28, bold=True, color=RGBColor.from_string(color_hex),
           align=PP_ALIGN.CENTER)
    _txbox(slide, left+0.1, top+0.62, 1.9, 0.35, label,
           size=9, color=WHITE, align=PP_ALIGN.CENTER)


def build_ppt(walks, m_walks, report_type: str, start_date: date, end_date: date):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    period_label = (str(start_date) if start_date == end_date
                    else f"{start_date.strftime('%b %d')} – {end_date.strftime('%b %d, %Y')}")
    total_barns  = sum(FARMS.values())
    total_checks = len(walks)
    total_mort   = sum(int(w.get("mortCount") or 0) for w in walks)
    total_loose  = sum(int(w.get("looseCount") or 0) for w in walks)
    total_flags  = sum(1 for w in walks if w.get("flags"))
    total_mw     = len(m_walks)
    avg_psi      = round(sum(float(w.get("waterPSI") or 0) for w in m_walks) / max(total_mw, 1), 1)
    avg_temp     = round(sum(float(w.get("temp") or 0) for w in m_walks) / max(total_mw, 1), 1)
    flagged_list = [(w.get("farm",""), w.get("house",""), "; ".join(w.get("flags",[])))
                    for w in walks if w.get("flags")]

    # ── Slide 1: Cover ─────────────────────────────────────────────────────────
    sl = _add_slide(prs)
    _rect(sl, 0, 0, 10, 7.5, "1a3a1a")
    _rect(sl, 0, 0, 10, 0.12, "4caf50")
    _rect(sl, 0, 7.38, 10, 0.12, "4caf50")
    _txbox(sl, 0.5, 1.5, 9, 1.2, "🏭 RUSHTOWN POULTRY",
           size=32, bold=True, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    _txbox(sl, 0.5, 2.7, 9, 0.8,
           f"{report_type.title()} Production Report",
           size=22, color=WHITE, align=PP_ALIGN.CENTER)
    _txbox(sl, 0.5, 3.5, 9, 0.6, period_label,
           size=16, color=RGBColor(0x7a, 0xb0, 0x7a), align=PP_ALIGN.CENTER)
    _txbox(sl, 0.5, 4.5, 9, 0.5,
           f"Generated {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
           size=11, color=GREY, align=PP_ALIGN.CENTER)

    # ── Slide 2: Executive Summary ─────────────────────────────────────────────
    sl = _add_slide(prs)
    _rect(sl, 0, 0, 10, 0.7, "1a3a1a")
    _txbox(sl, 0.3, 0.1, 9.4, 0.5, "Executive Summary",
           size=18, bold=True, color=WHITE)

    kpis = [
        ("Barn Checks", total_checks, "4caf50" if total_checks >= total_barns else "d69e2e"),
        ("Mortality", total_mort, "e53e3e" if total_mort > 0 else "4caf50"),
        ("Loose Birds", total_loose, "e53e3e" if total_loose > 0 else "4caf50"),
        ("Flagged Barns", total_flags, "e53e3e" if total_flags > 0 else "4caf50"),
    ]
    for i, (lbl, val, clr) in enumerate(kpis):
        _kpi_box(sl, 0.4 + i * 2.3, 1.0, lbl, val, clr)

    kpis2 = [
        ("Morning Walks", total_mw, "4caf50"),
        ("Avg Water PSI", avg_psi, "4caf50" if 10 <= avg_psi <= 60 else "e53e3e"),
        ("Avg Temp °F", avg_temp, "4caf50"),
    ]
    for i, (lbl, val, clr) in enumerate(kpis2):
        _kpi_box(sl, 1.75 + i * 2.3, 2.4, lbl, val, clr)

    # Summary text
    checks_pct = round(total_checks / total_barns * 100) if total_barns else 0
    summary_lines = [
        f"• {total_checks}/{total_barns} barns checked ({checks_pct}% completion)",
        f"• {total_mort} total mortality recorded" + (" ⚠" if total_mort > 0 else " ✓"),
        f"• {total_flags} barn checks flagged issues" + (" — Work Orders created" if total_flags > 0 else " — All clear"),
        f"• {total_mw} morning walks completed by Lead/WNO",
        f"• Average water pressure: {avg_psi} PSI  |  Average temp: {avg_temp}°F",
    ]
    _txbox(sl, 0.4, 3.7, 9.2, 2.8, "\n".join(summary_lines),
           size=13, color=RGBColor(0x2a, 0x2a, 0x2a))

    # ── Slide 3: Barn Check Status ─────────────────────────────────────────────
    sl = _add_slide(prs)
    _rect(sl, 0, 0, 10, 0.7, "1a3a1a")
    _txbox(sl, 0.3, 0.1, 9.4, 0.5, "Barn Check Status",
           size=18, bold=True, color=WHITE)

    y = 0.85
    for farm, cnt in FARMS.items():
        _txbox(sl, 0.3, y, 9.4, 0.35, f"📍 {farm}",
               size=12, bold=True, color=GREEN_LIGHT)
        y += 0.35
        farm_walks = {w.get("house"): w for w in walks if w.get("farm") == farm}
        cols = 4
        for i in range(1, cnt + 1):
            col = (i - 1) % cols
            row = (i - 1) // cols
            w = farm_walks.get(str(i))
            flagged = w and bool(w.get("flags"))
            done    = bool(w)
            fill    = "2a1a1a" if flagged else ("1a3a1a" if done else "2a2a2a")
            icon    = "⚠" if flagged else ("✓" if done else "—")
            lx = 0.3 + col * 2.4
            ly = y + row * 0.75
            _rect(sl, lx, ly, 2.2, 0.65, fill)
            mort_str = f"  💀{w['mortCount']}" if (w and w.get("mortCount")) else ""
            _txbox(sl, lx + 0.05, ly + 0.05, 2.1, 0.55,
                   f"Barn {i}  {icon}{mort_str}",
                   size=11, bold=True,
                   color=(RED if flagged else (GREEN_LIGHT if done else GREY)),
                   align=PP_ALIGN.CENTER)
        rows_used = ((cnt - 1) // cols) + 1
        y += rows_used * 0.75 + 0.1

    # ── Slide 4: Flags & Issues ────────────────────────────────────────────────
    sl = _add_slide(prs)
    _rect(sl, 0, 0, 10, 0.7, "2a0f0f")
    _txbox(sl, 0.3, 0.1, 9.4, 0.5, f"⚠ Flags & Issues  ({len(flagged_list)} barns)",
           size=18, bold=True, color=WHITE)

    if not flagged_list:
        _txbox(sl, 0.5, 1.5, 9, 1, "✓ No flags recorded — all barns clear",
               size=18, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    else:
        y = 0.85
        for farm_n, house, flags_str in flagged_list[:12]:
            _rect(sl, 0.3, y, 9.4, 0.6, "2a1010")
            _txbox(sl, 0.4, y + 0.05, 2.5, 0.5,
                   f"{farm_n} — Barn {house}", size=12, bold=True, color=RED)
            _txbox(sl, 2.9, y + 0.05, 6.7, 0.5,
                   flags_str, size=10, color=RGBColor(0xf0, 0x90, 0x90))
            y += 0.65
        if len(flagged_list) > 12:
            _txbox(sl, 0.3, y, 9, 0.4,
                   f"+ {len(flagged_list)-12} more — see Excel report",
                   size=10, color=GREY)

    # ── Slide 5: Morning Walk Summary ─────────────────────────────────────────
    sl = _add_slide(prs)
    _rect(sl, 0, 0, 10, 0.7, "0d1f3a")
    _txbox(sl, 0.3, 0.1, 9.4, 0.5, "☀ Morning Walk Summary",
           size=18, bold=True, color=WHITE)

    if not m_walks:
        _txbox(sl, 0.5, 1.5, 9, 1, "No morning walks recorded for this period.",
               size=16, color=GREY, align=PP_ALIGN.CENTER)
    else:
        _txbox(sl, 0.3, 0.75, 9.4, 0.3,
               f"Lead / WNO completions: {total_mw}/{total_barns}   |   Avg PSI: {avg_psi}   |   Avg Temp: {avg_temp}°F",
               size=11, color=RGBColor(0x6a, 0x90, 0xd9))
        y = 1.1
        mw_by_key = {(w.get("farm"), w.get("house")): w for w in m_walks}
        for farm, cnt in FARMS.items():
            _txbox(sl, 0.3, y, 9, 0.3, f"📍 {farm}",
                   size=11, bold=True, color=RGBColor(0x6a, 0x90, 0xd9))
            y += 0.3
            for i in range(1, cnt + 1):
                mw = mw_by_key.get((farm, str(i)))
                if mw:
                    ee_str = f"  EE: {mw['eeCount']}" if mw.get("eeCount") is not None else ""
                    flag_str = "  ⚠ " + "; ".join(mw["flags"]) if mw.get("flags") else ""
                    line = (f"  Barn {i}: {mw.get('employee','')}  |  "
                            f"{mw.get('waterPSI','')} PSI  |  {mw.get('temp','')}°F"
                            f"{ee_str}{flag_str}")
                    clr = RED if mw.get("flags") else RGBColor(0xd0, 0xe8, 0xd0)
                    _txbox(sl, 0.3, y, 9.4, 0.32, line, size=10, color=clr)
                    y += 0.32
                    if y > 7.0:
                        break
            if y > 7.0:
                break

    # ── Weekly/Monthly extra slides ────────────────────────────────────────────
    if report_type in ("weekly", "monthly"):
        # Per-day breakdown
        all_dates = sorted(set(w.get("date","") for w in walks if w.get("date")))
        if all_dates:
            sl = _add_slide(prs)
            _rect(sl, 0, 0, 10, 0.7, "1a3a1a")
            _txbox(sl, 0.3, 0.1, 9.4, 0.5, "Daily Breakdown",
                   size=18, bold=True, color=WHITE)
            headers = ["Date", "Checks", "Mortality", "Loose", "Flagged", "AM Walks"]
            col_w = [1.6, 1.1, 1.2, 1.1, 1.2, 1.2]
            x_starts = [0.3]
            for w in col_w[:-1]:
                x_starts.append(x_starts[-1] + w)
            y = 0.75
            for ci, (h, cw, cx) in enumerate(zip(headers, col_w, x_starts)):
                _rect(sl, cx, y, cw - 0.05, 0.35, "1a3a1a")
                _txbox(sl, cx, y, cw, 0.35, h, size=10, bold=True,
                       color=WHITE, align=PP_ALIGN.CENTER)
            y += 0.38
            for di, d_str in enumerate(all_dates):
                day_w  = [w for w in walks   if w.get("date") == d_str]
                day_mw = [w for w in m_walks if w.get("date") == d_str]
                row_data = [
                    d_str,
                    str(len(day_w)),
                    str(sum(int(w.get("mortCount") or 0) for w in day_w)),
                    str(sum(int(w.get("looseCount") or 0) for w in day_w)),
                    str(sum(1 for w in day_w if w.get("flags"))),
                    str(len(day_mw)),
                ]
                fill = "1a2a1a" if di % 2 == 0 else "121a12"
                for ci, (val, cw, cx) in enumerate(zip(row_data, col_w, x_starts)):
                    _rect(sl, cx, y, cw - 0.05, 0.32, fill)
                    clr = RED if (ci == 4 and val != "0") else (
                          RED if (ci == 2 and val != "0") else WHITE)
                    _txbox(sl, cx, y, cw, 0.32, val, size=10,
                           color=clr, align=PP_ALIGN.CENTER)
                y += 0.33
                if y > 7.2:
                    break

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Email ──────────────────────────────────────────────────────────────────────
def send_email(api_key: str, report_type: str, period: str,
               excel_bytes: bytes, ppt_bytes: bytes):
    resend.api_key = api_key
    subject = f"Rushtown Poultry — {report_type.title()} Report ({period})"
    body = (f"<p>Good morning,</p>"
            f"<p>Please find attached the <strong>{report_type}</strong> production report "
            f"for <strong>{period}</strong>.</p>"
            f"<ul>"
            f"<li>📊 Excel spreadsheet with full barn check and morning walk data</li>"
            f"<li>📑 PowerPoint presentation with summary and flagged items</li>"
            f"</ul>"
            f"<p>This report was generated automatically by the Rushtown Poultry Operations Hub.</p>")

    fname = period.replace(" ", "_").replace(",", "").replace("–", "to")

    params = {
        "from": "Rushtown Poultry <onboarding@resend.dev>",
        "to": [SEND_TO],
        "subject": subject,
        "html": body,
        "attachments": [
            {
                "filename": f"Rushtown_{report_type}_{fname}.xlsx",
                "content": list(excel_bytes),
            },
            {
                "filename": f"Rushtown_{report_type}_{fname}.pptx",
                "content": list(ppt_bytes),
            },
        ],
    }
    resp = resend.Emails.send(params)
    print(f"Email sent: {resp}")


# ── Main ───────────────────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--type", choices=["daily", "weekly", "monthly"], required=True)
    args = parser.parse_args()

    sg_key = os.environ.get("RESEND_API_KEY")
    if not sg_key:
        sys.exit("Missing RESEND_API_KEY env var")

    db = init_firebase()
    start_d, end_d = date_range(args.type)
    print(f"Generating {args.type} report: {start_d} → {end_d}")

    walks   = fetch_walks(db, start_d, end_d)
    m_walks = fetch_morning_walks(db, start_d, end_d)
    print(f"  Barn checks: {len(walks)}, Morning walks: {len(m_walks)}")

    period = str(start_d) if start_d == end_d else f"{start_d} to {end_d}"

    # EOS snapshot first — most important output, write it before anything that can fail
    try:
        snapshot = build_eos_snapshot(db, date.today())
        write_eos_snapshot(snapshot, path="data/eos-snapshot.json")
    except Exception as e:
        print(f"WARNING — failed to write EOS snapshot: {e}", file=sys.stderr)

    excel  = build_excel(walks, m_walks, args.type, start_d, end_d)
    ppt    = build_ppt(walks, m_walks, args.type, start_d, end_d)

    # Email — best effort. If Resend rejects (test-mode, unverified domain, etc),
    # log it and keep going so the snapshot still gets committed by the workflow.
    try:
        send_email(sg_key, args.type, period, excel, ppt)
    except Exception as e:
        print(f"WARNING — email send failed (continuing anyway): {e}", file=sys.stderr)

    print("Done.")


if __name__ == "__main__":
    main()
